import requests
from pptx import Presentation

# 1. Set up your API request
url = "https://rattle-rm-reporting-hub.prod.crto.in/pivot_report/api/paid"

headers = {
    "accept": "*/*",
    "content-type": "application/json",
    "origin": "https://rattle-rm-reporting-hub.prod.crto.in",
    "referer": "https://rattle-rm-reporting-hub.prod.crto.in/pivot_report/pivot/",
    "cookie": "_ga=GA1.1.604358013.1708096317; _hjSessionUser_3428851=eyJpZCI6IjhlNTFlOTcwLWE5NjctNTQ3My05N2UxLTAzOTA2NGVlZWM3ZCIsImNyZWF0ZWQiOjE3MzgwMTg2MDc3NjIsImV4aXN0aW5nIjp0cnVlfQ==; ajs_user_id=425424; ajs_anonymous_id=f38bf151-f8e2-40c1-8f77-d7704837fe1d; amplitude_idundefinedcrto.in=eyJvcHRPdXQiOmZhbHNlLCJzZXNzaW9uSWQiOm51bGwsImxhc3RFdmVudFRpbWUiOm51bGwsImV2ZW50SWQiOjAsImlkZW50aWZ5SWQiOjAsInNlcXVlbmNlTnVtYmVyIjowfQ==; amplitude_id_339b32621c58a6ab1ad15e35ec90742fcrto.in=eyJkZXZpY2VJZCI6Ijg5ZGJhYWNiLTA4NDUtNGFjNC1iOWFjLTUzYzkzZWYxNGM2YlIiLCJ1c2VySWQiOiI0MjU0MjQiLCJvcHRPdXQiOmZhbHNlLCJzZXNzaW9uSWQiOjE3NDM2MjIxMDYyMDgsImxhc3RFdmVudFRpbWUiOjE3NDM2MjIxMTI0OTIsImV2ZW50SWQiOjQyLCJpZGVudGlmeUlkIjoxMSwic2VxdWVuY2VOdW1iZXIiOjUzfQ==; _ga_YJMCDYV8GD=GS1.1.1746036937.291.0.1746036937.0.0.0; session=.eJytVF1v0zAU_StRnvuVrO1g2hBp2m1llHaNCmyaFHnOTeLWsYPttA1o_53rjLEOifHCk31Prs899ys_XLkxJK40qJiJVLonP1woCOPuibvugNAc6vdUMQOyQ2XhttyUKW1iQQpAlw-EbrQUCHPyjE6aZwi-9HJ-41WhY5a4J31_0Pf7-FhmTBxEdB9aj7qoTCDegmIpA4UOOz-HwagO9mMRfSyEPJ_UV3Q2zhbB2-PLLyVNBzcflnSnjuX11c1uPRzUKoqre8_ofm--n-1vp347jb_GX5Yr0x8uhtl9LjPP87NNNDoaJuLiVpe7ILspw8n18PO1dzH_drUrd5fnl8F2f_l9jdqV5DanYIJ3KEoua4DXk8GaVZzHfysG0zFJCvskJVxD6x_1Z4IZRnickBp9_J7vt3vDtu8dyuGwBWR40-m1XCorYZT1XUVWPmRMWnnBbLKchoHFEiiJMgUIg_hCsS0x4MyI2oBxIsJBoyFIBo1Hyy3IPk50-RRFYHIttwRVMK2RW9sZ6iNTBDx1IlBbcCY2JWcJpVSWYWDjzJ7tIdrL2cJZEMpEhsBbBEYgaF6gCrS9HgLjaOGMic7vJVGJBT0b5E_Qt7nRJmsn4KDMI6N3ZCnno7m9W3VToVmWG23tRg_bSvMsyWs0gUHdoJyAGrZlpnaiqkBJtTNJWOP15iDaZwY7i1nxISllZZzVFAHfitdya6daxE13Y9XEwVL5FiO_dMZGSo6gh21rev5yMV-di8Q2Nfx09vS1NV-dnbNihQS6NQ7PKpKUSib2-nRSZaQ9mV1gKv6yqdOD2a7soL9Y04O-xzQHugH0MKoCO2qEmpgVvzJ1TzUKLo2TQArq7M5tzjvXMXUJaBYyqTigrRVFs6sNMYx2G5Yu0RqM7jKRwL4drvPrY3PU7qz1nfvutPvI--5_8XdAt8Pl8GO1Ow7_f4iywr9ZbWOQdkjq78HVnzE4ExsnV5C-XoVxdDutYHTRoRrfOwo4-mtT477mAAYp3Yefkc3y2w.aBJoyg.MdHjFzErt3TyRzbuZqYyBjA7g2Y"
}

payload = {
    "retailers": ["906"],
    "accounts": [],
    "private-market": [],
    "countries": [],
    "dimensions": ["retailer"],
    "start_date": "2025-03-01",
    "end_date": "2025-03-31",
    "columns": "",
    "rsxRmp": ["RMP"],
    "metrics": ["impressions", "clicks", "spend"],
    "timezone": "UTC",
    "currency": "Default",
    "rsxCampaigns": [],
    "rmpCampaigns": [],
    "lineItemFilter": [],
    "globalAccountFilter": [],
    "campaignType": [1, 2],
    "saleType": ["online", "offline"],
    "attrWindow": "campaign_default",
    "attrLevelClick": "campaign_default",
    "attrLevelImp": "campaign_default",
    "attrRsxLevel": ["Default"],
    "environment": [],
    "pageType": [],
    "topKeywordFilter": None,
    "topKeywordMetric": "impressions",
    "keywordFilter": [],
    "sspNetwork": [],
    "skuType": "retailer_sku_key",
    "skuFilter": [],
    "productCategoryLike": "",
    "productCategoryLevel": 1,
    "retailerProductCategory": [],
    "productCategory": [],
    "parentCategoryAgg": True,
    "isLockout": "",
    "isManaged": "",
    "optOut": "",
    "agency": [],
    "accountStrategist": [],
    "salesRepresentative": [],
    "campaignManager": [],
    "download": False,
    "query": False
}

# 2. Make the request
response = requests.post(url, headers=headers, json=payload)

if response.status_code != 200:
    print(f"Request failed with status {response.status_code}")
    print(response.text)
    exit()

# 3. Parse the JSON response
api_data = response.json()
row = api_data.get("data", [{}])[0]

retailer_name = row.get("retailer_name", "Unknown")
impressions = row.get("impressions", 0)
clicks = row.get("clicks", 0)
spend = row.get("spend", 0.0)
currency = row.get("currency", "USD")

# 4. Create a dictionary of placeholders and replacements
replacements = {
    "{{name}}": retailer_name,
    "{{impressions}}": f"{impressions:,}",
    "{{clicks}}": f"{clicks:,}",
    "{{spend}}": f"{currency} {spend:,.2f}"
}

# 5. Open PowerPoint and replace text
prs = Presentation("template.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for key, value in replacements.items():
                if key in shape.text:
                    shape.text = shape.text.replace(key, value)

# 6. Save the result
prs.save("output.pptx")
print("Presentation updated successfully.")
