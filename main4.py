import requests
from pptx import Presentation


#####

retailerID = "906" #albertsons
campaignID = [214650]
lineItemsPivot = [685422]
lineItemsCapout = "685422"
brandName = "Nutpods"
presentationName = "Nutpods Snapshot Deck (May 2025).pptx"

#####


# Set up your API details for the Pivot API
urlPivot = "https://rattle-rm-reporting-hub.prod.crto.in/pivot_report/api/paid"
headersPivot = {
    "accept": "*/*",
    "content-type": "application/json",
    "origin": "https://rattle-rm-reporting-hub.prod.crto.in",
    "referer": "https://rattle-rm-reporting-hub.prod.crto.in/pivot_report/pivot/",
    "cookie": "_ga=GA1.1.604358013.1708096317; _hjSessionUser_3428851=eyJpZCI6IjhlNTFlOTcwLWE5NjctNTQ3My05N2UxLTAzOTA2NGVlZWM3ZCIsImNyZWF0ZWQiOjE3MzgwMTg2MDc3NjIsImV4aXN0aW5nIjp0cnVlfQ==; ajs_user_id=425424; ajs_anonymous_id=f38bf151-f8e2-40c1-8f77-d7704837fe1d; amplitude_idundefinedcrto.in=eyJvcHRPdXQiOmZhbHNlLCJzZXNzaW9uSWQiOm51bGwsImxhc3RFdmVudFRpbWUiOm51bGwsImV2ZW50SWQiOjAsImlkZW50aWZ5SWQiOjAsInNlcXVlbmNlTnVtYmVyIjowfQ==; amplitude_id_339b32621c58a6ab1ad15e35ec90742fcrto.in=eyJkZXZpY2VJZCI6Ijg5ZGJhYWNiLTA4NDUtNGFjNC1iOWFjLTUzYzkzZWYxNGM2YlIiLCJ1c2VySWQiOiI0MjU0MjQiLCJvcHRPdXQiOmZhbHNlLCJzZXNzaW9uSWQiOjE3NDM2MjIxMDYyMDgsImxhc3RFdmVudFRpbWUiOjE3NDM2MjIxMTI0OTIsImV2ZW50SWQiOjQyLCJpZGVudGlmeUlkIjoxMSwic2VxdWVuY2VOdW1iZXIiOjUzfQ==; _ga_YJMCDYV8GD=GS1.1.1746036937.291.0.1746036937.0.0.0; session=.eJytVF1v0zAU_StRnvuVrO1g2hBp2m1llHaNCmyaFHnOTeLWsYPttA1o_53rjLEOifHCk31Prs899ys_XLkxJK40qJiJVLonP1woCOPuibvugNAc6vdUMQOyQ2XhttyUKW1iQQpAlw-EbrQUCHPyjE6aZwi-9HJ-41WhY5a4J31_0Pf7-FhmTBxEdB9aj7qoTCDegmIpA4UOOz-HwagO9mMRfSyEPJ_UV3Q2zhbB2-PLLyVNBzcflnSnjuX11c1uPRzUKoqre8_ofm--n-1vp347jb_GX5Yr0x8uhtl9LjPP87NNNDoaJuLiVpe7ILspw8n18PO1dzH_drUrd5fnl8F2f_l9jdqV5DanYIJ3KEoua4DXk8GaVZzHfysG0zFJCvskJVxD6x_1Z4IZRnickBp9_J7vt3vDtu8dyuGwBWR40-m1XCorYZT1XUVWPmRMWnnBbLKchoHFEiiJMgUIg_hCsS0x4MyI2oBxIsJBoyFIBo1Hyy3IPk50-RRFYHIttwRVMK2RW9sZ6iNTBDx1IlBbcCY2JWcJpVSWYWDjzJ7tIdrL2cJZEMpEhsBbBEYgaF6gCrS9HgLjaOGMic7vJVGJBT0b5E_Qt7nRJmsn4KDMI6N3ZCnno7m9W3VToVmWG23tRg_bSvMsyWs0gUHdoJyAGrZlpnaiqkBJtTNJWOP15iDaZwY7i1nxISllZZzVFAHfitdya6daxE13Y9XEwVL5FiO_dMZGSo6gh21rev5yMV-di8Q2Nfx09vS1NV-dnbNihQS6NQ7PKpKUSib2-nRSZaQ9mV1gKv6yqdOD2a7soL9Y04O-xzQHugH0MKoCO2qEmpgVvzJ1TzUKLo2TQArq7M5tzjvXMXUJaBYyqTigrRVFs6sNMYx2G5Yu0RqM7jKRwL4drvPrY3PU7qz1nfvutPvI--5_8XdAt8Pl8GO1Ow7_f4iywr9ZbWOQdkjq78HVnzE4ExsnV5C-XoVxdDutYHTRoRrfOwo4-mtT477mAAYp3Yefkc3y2w.aBJoyg.MdHjFzErt3TyRzbuZqYyBjA7g2Y"
}

urlCapout = "https://rattle-rm-reporting-hub.prod.crto.in/capout/api/v2/table"
headersCapout = {
    "accept": "*/*",
    "referer": "https://rattle-rm-reporting-hub.prod.crto.in/capout/reports/10146",
    "user-agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/135.0.0.0 Safari/537.36",
    "cookie": "_ga=GA1.1.604358013.1708096317; _hjSessionUser_3428851=eyJpZCI6IjhlNTFlOTcwLWE5NjctNTQ3My05N2UxLTAzOTA2NGVlZWM3ZCIsImNyZWF0ZWQiOjE3MzgwMTg2MDc3NjIsImV4aXN0aW5nIjp0cnVlfQ==; ajs_user_id=425424; ajs_anonymous_id=f38bf151-f8e2-40c1-8f77-d7704837fe1d; amplitude_idundefinedcrto.in=eyJvcHRPdXQiOmZhbHNlLCJzZXNzaW9uSWQiOm51bGwsImxhc3RFdmVudFRpbWUiOm51bGwsImV2ZW50SWQiOjAsImlkZW50aWZ5SWQiOjAsInNlcXVlbmNlTnVtYmVyIjowfQ==; amplitude_id_339b32621c58a6ab1ad15e35ec90742fcrto.in=eyJkZXZpY2VJZCI6Ijg5ZGJhYWNiLTA4NDUtNGFjNC1iOWFjLTUzYzkzZWYxNGM2YlIiLCJ1c2VySWQiOiI0MjU0MjQiLCJvcHRPdXQiOmZhbHNlLCJzZXNzaW9uSWQiOjE3NDM2MjIxMDYyMDgsImxhc3RFdmVudFRpbWUiOjE3NDM2MjIxMTI0OTIsImV2ZW50SWQiOjQyLCJpZGVudGlmeUlkIjoxMSwic2VxdWVuY2VOdW1iZXIiOjUzfQ==; session=.eJytVNtS2zAQ_RWPnnOzm6SUKZ06F0hK04R40haGGY-Q146ILLmSnMTt8O9dGSiBGdqXPkl7vDp79uZfRG0sjUsDOuYyVeT4F4GcckGOyW0LpBFQfWSaW1AtpnLSICnXxsaS5oAunyjbGCURFvQJHdfPEHzu5f3By9zEPCHH3aDXDbr4WGVcHkQkd417XUwlEG9B85SDRoddsIbeoAr3Ixl9zqU6HVfnbDbKFuG7t5NvBUt7l5-WbKffqovzy91tv1fpKC5vfGu6nfl-tr-aBs00_h5_W65st7_oZzdrlfl-kG2iwZt-Is-uTLELs8tiOL7of73wz-Y_znfFbnI6Cbf7yc9b1K6VcDmFY7xDXghVAfw9GaxZKUT8WjG4iWmSuycpFQYa_6g_l9xyKuKEVugTdIKg2ek3A_9QjoAtIMNRq9MgTJXSaue7ipx8yLhy8sLZeDkdhg5LoKDa5iAt4gvNt9SCN6N6A9aLqACDhqQZ1B4NktN9nJjiMYrE5BqkAJ1zY5DbuBnqIlMEIvUi0Fvwxi4lbwmF0o6h5-LMnuw-2svZwltQxmWGwDsEBiDZOkcVaPsdBEbRwhtRs75RVCcO9F2Ql2DgcmN11l4oQNt7Rv-No5wP5u7u1E2l4dnaGmfXevhW2SdJfq0JLOoG7YXM8i23lReVOUqqvHHCa6-jg2hfOewc5sQPaaFK662mCAROvFFbN9Uyrrsb6zoOlipwGH3QGVulBII-tq3u-fPF_OtcJK6pwy8nj18b89XJKc9XSGAao-FJSZNCq8RdH0-mrXIndwvM5CubOj2Y7dIN-rM1Peh7zNbANoAeVpfgRo0yG_P8IVPy3qDgwnoJpKBPrkl9XhPPVgWgmaukFIC20QzNtrHUctauWdrUGLCmzWUC--aI_ehdNINx69Zckw_v2_e8H_4XfwtMU345upqsbor_H6Io8W9WuRi0OaTVz_D8ZQzB5cZba0j_UYXoalrC4KzFDL73NAj0N7bCfV0DWKQkd78BjHDy7A.aBO8uA.HkZdannJVmoVewTgx3hy8jYNB-I; _ga_YJMCDYV8GD=GS1.1.1746123959.294.0.1746123960.0.0.0"
}


# Template for Pivot payload
def get_payload_pivot(start_date, end_date, line_item_filter=None, by_keyword=None, by_category=None):
    dimensions = ["retailer", "campaign"]  # Default dimensions
    
    if line_item_filter:
        dimensions.append("lineItem")  # Add lineItem if we're using line-item filters

    if by_keyword:
        dimensions.append("keyword")  # Add keyword if we're using keyword filtering

    if by_category:
        dimensions.append("breadcrumb")  # Add keyword if we're using keyword filtering
    
    payloadPivot = {
        "retailers": [retailerID],
        "accounts": [],
        "private-market": [],
        "countries": [],
        "dimensions": dimensions,
        "columns": "",
        "rsxRmp": ["RMP"],
        "metrics": ["impressions", "clicks", "spend", "sales", "units", "CPC Demand", "ROAS Demand"],
        "timezone": "UTC",
        "currency": "Default",
        "rsxCampaigns": [],
        "rmpCampaigns": campaignID,
        "globalAccountFilter": [],
        "campaignType": [1, 2],
        "saleType": ["online", "offline"],
        "attrWindow": "campaign_default",
        "attrLevelClick": "campaign_default",
        "attrLevelImp": "campaign_default",
        "attrRsxLevel": ["Default"],
        "environment": [],
        "pageType": [],
        "topKeywordFilter": 16,
        "topKeywordMetric": "spend",
        "keywordFilter": [],
        "sspNetwork": [],
        "skuType": "retailer_sku_key",
        "skuFilter": [],
        "productCategoryLike": "",
        "productCategoryLevel": 1,
        "retailerProductCategory": [],
        "productCategory": [],
        "parentCategoryAgg": True,
        "isLockout": "",
        "isManaged": "",
        "optOut": "",
        "agency": [],
        "accountStrategist": [],
        "salesRepresentative": [],
        "campaignManager": [],
        "download": False,
        "query": False,
        "start_date": start_date,
        "end_date": end_date,
    }
    
    #gives line-item data if true.
    if line_item_filter:
        payloadPivot["lineItemFilter"] = line_item_filter
    else:
        payloadPivot["lineItemFilter"] = []  # Leave empty to get campaign-level data

    #gives keyword data if true.
    if by_keyword:
        payloadPivot["byKeyword"] = by_keyword
    else:
        payloadPivot["byKeyword"] = []  # Leave empty to get campaign-level data

    #gives L3 category data if true.
    if by_category:
        payloadPivot["byCategory"] = by_category
    else:
        payloadPivot["byCategory"] = []  # Leave empty to get campaign-level data

    return payloadPivot

def get_payload_capout(start_date, end_date, line_item_filter_capout=None):
    dims = ["campaign_id"]
    if line_item_filter_capout:
        dims.append("line_item_id")

    p = {
        "retailer_id":    retailerID,
        "campaign_id":    campaignID,
        "current_period_start": start_date,
        "current_period_end":   end_date,
        "metrics":        "alert_hour",
        "dimensions":     ",".join(dims),
        "timezone":       "UTC",
        "download":       "false",
        "currency_code":  "USD",
    }

    if line_item_filter_capout:
        p["line_item_id"] = (
            line_item_filter_capout 
            if isinstance(line_item_filter_capout, str) 
            else ",".join(map(str, line_item_filter_capout))
        )

    return p

# Call the Pivot API
def get_data_pivot(start_date, end_date, line_item_filter=None, by_keyword=None, by_category=None):

    #Pivot specific API response
    payloadPivot = get_payload_pivot(start_date, end_date, line_item_filter, by_keyword, by_category)
    responsePivot = requests.post(urlPivot, headers=headersPivot, json=payloadPivot)
    if responsePivot.status_code != 200:
        print(f"Request failed: {responsePivot.status_code}")
        print(responsePivot.text)
        return []
    return responsePivot.json().get("data", [])

def get_data_capout(start_date, end_date, line_item_filter_capout=None):

    #Pivot specific API response

    payloadCapout = get_payload_capout(start_date, end_date, line_item_filter_capout)

    responseCapout = requests.get(urlCapout, headers=headersCapout, params=payloadCapout, timeout=20)
    
    if responseCapout.status_code != 200:
        print(f"Request failed: {responseCapout.status_code}")
        print(responseCapout.text)
        return []
    
    return responseCapout.json().get("data", [])

#convert pivot time to capout time
from datetime import datetime
def to_capout_iso(date_str):
    # parse "YYYY-MM-DD" then format as "YYYY-MM-DDT00:00:00"
    dt = datetime.strptime(date_str, "%Y-%m-%d")
    return dt.strftime("%Y-%m-%dT00:00:00")


# Replace text in shapes
def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for key, value in replacements.items():
                if key in run.text:
                    run.text = run.text.replace(key, value)

"""
#Replace text in table
def replace_text_in_table(table, replacements):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    print(f"RUN: {repr(run.text)}")
                    for key, value in replacements.items():
                        if key in run.text:
                            run.text = run.text.replace(key, value)
"""

from pptx.dml.color import RGBColor
def replace_text_in_table(table, replacements):
    for row in table.rows:
        for cell in row.cells:
            tf = cell.text_frame
            for p in list(tf.paragraphs):
                text = p.text
                if not any(key in text for key in replacements):
                    continue

                # 1) capture formatting from the first run
                first = p.runs[0]
                fnt = first.font
                fmt = {
                    "name":   fnt.name,
                    "size":   fnt.size,
                    "bold":   fnt.bold,
                    "italic": fnt.italic,
                    "color":  (fnt.color.rgb if fnt.color.type is not None else None),
                }

                # 2) do the replacements on the full paragraph text
                full = text
                for key, val in replacements.items():
                    s = "" if val is None else str(val)
                    full = full.replace(key, s)

                # 3) remove all existing runs via the paragraph XML
                for run in list(p.runs):
                    p._p.remove(run._r)

                # 4) add a single new run with the replaced text
                new_run = p.add_run()
                new_run.text = full

                # 5) re-apply the captured formatting
                nf = new_run.font
                nf.name   = fmt["name"]
                nf.size   = fmt["size"]
                nf.bold   = fmt["bold"]
                nf.italic = fmt["italic"]
                if fmt["color"] is not None:
                    nf.color.rgb = fmt["color"]


            
# Fill a slide with data (either campaign-level, line-item, keyword, or category)
def fill_slide(prs, slide_index, start_date, end_date, use_line_items=False, by_keyword=False, by_category=False):
    
    #write internal line-item IDs here
    line_item_filter = lineItemsPivot if use_line_items else None
    line_item_filter_capout = lineItemsCapout if use_line_items else None

    #pivot 
    pivot_start = start_date
    pivot_end   = end_date

    # capout wants "YYYY-MM-DDT00:00:00"
    capout_start = to_capout_iso(start_date)
    capout_end   = to_capout_iso(end_date)
    
    dataPivot = get_data_pivot(pivot_start, pivot_end, line_item_filter, by_keyword, by_category)
    dataCapout = get_data_capout(capout_start, capout_end, line_item_filter_capout)

    snapshot_date = human_date(start_date) + " - " + human_date(end_date)
    
    default_row = dataPivot[0] if dataPivot else {}
    capout_row = dataCapout[0] if dataCapout else {}

    replacements = {
        "{{name}}": default_row.get("retailer_name", "Unknown"),
        "{{impressions}}": f"{default_row.get('impressions', 0):,}",
        "{{clicks}}": f"{default_row.get('clicks', 0):,}",
        "{{spend}}": f"${int(default_row.get('spend', 0)):,}",
        "{{sales}}": f"${int(default_row.get('sales', 0)):,}",
        "{{units}}": f"{default_row.get('units', 0):,}",
        "{{ROAS}}": f"${default_row.get('ROAS Demand', 0):,.2f}",
        "{{cap out}}": format_hour_12h(capout_row.get("alert_hour_current", 0)) + " EST",
        "{{CPC}}": f"${default_row.get('CPC Demand', 0):,.2f}",
        "{{snapshot date}}": snapshot_date,
        "{{keyword}}": default_row.get("keyword"),

        #ENTER YOURSELF
        "{{BRAND NAME}}": brandName,
    }

    if use_line_items:
        dataPivot.sort(key=lambda x: x["spend"], reverse=True)


        #pick up here. you just added the max_items line. test if this works on pilgrim's, and then a campaign that has 3 line-items.


        
        max_items = len(line_item_filter)
        for idx, line_item_row in enumerate(dataPivot[:max_items], start=1):
            replacements[f"{{{{impressions_lineitem_{idx}}}}}"] = f"{line_item_row.get('impressions', 0):,}"
            replacements[f"{{{{clicks_lineitem_{idx}}}}}"] = f"{line_item_row.get('clicks', 0):,}"
            replacements[f"{{{{spend_lineitem_{idx}}}}}"] = f"${int(line_item_row.get('spend', 0)):,}"
            replacements[f"{{{{sales_lineitem_{idx}}}}}"] = f"${int(line_item_row.get('sales', 0)):,}"
            replacements[f"{{{{units_lineitem_{idx}}}}}"] = f"{line_item_row.get('units', 0):,}"
            replacements[f"{{{{ROAS_lineitem_{idx}}}}}"] = f"${line_item_row.get('ROAS Demand', 0):,.2f}"
            replacements[f"{{{{CPC_lineitem_{idx}}}}}"] = f"${line_item_row.get('CPC Demand', 0):,.2f}"
            replacements[f"{{{{name_lineitem_{idx}}}}}"] = line_item_row.get("line_item_name", f"LineItem {idx}")
        for idx, line_item_row in enumerate(dataCapout[:max_items], start=1):
            # grab the alert hour, default to 0, cast to int
            ah = int(line_item_row.get("alert_hour_current", 0))
            
            # format with commas and append ":00 EST" if you need the same style
            replacements[f"{{{{cap_lineitem_{idx}}}}}"] = format_hour_12h(ah) + " EST"
            

    if by_keyword:
        for idx, keywordRow in enumerate(dataPivot[:15], start=1):
            replacements[f"{{{{keyword_{idx}}}}}"] = keywordRow.get("keyword", f"Keyword {idx}")
            replacements[f"{{{{impressions_keyword_{idx}}}}}"] = f"{keywordRow.get('impressions', 0):,}"
            replacements[f"{{{{clicks_keyword_{idx}}}}}"] = f"{keywordRow.get('clicks', 0):,}"
            replacements[f"{{{{spend_keyword_{idx}}}}}"] = f"${int(keywordRow.get('spend', 0)):,}"
            replacements[f"{{{{sales_keyword_{idx}}}}}"] = f"${int(keywordRow.get('sales', 0)):,}"
            replacements[f"{{{{units_keyword_{idx}}}}}"] = f"{keywordRow.get('units', 0):,}"
            replacements[f"{{{{ROAS_keyword_{idx}}}}}"] = f"${keywordRow.get('ROAS Demand', 0):,.2f}"
            replacements[f"{{{{CPC_keyword_{idx}}}}}"] = f"${keywordRow.get('CPC Demand', 0):,.2f}"

    if by_category:
        dataPivot.sort(key=lambda x: x["spend"], reverse=True)
        for idx, categoryRow in enumerate(dataPivot[:15], start=1):
            replacements[f"{{{{category_{idx}}}}}"] = categoryRow.get("breadcrumb", f"Category {idx}")
            replacements[f"{{{{impressions_category_{idx}}}}}"] = f"{categoryRow.get('impressions', 0):,}"
            replacements[f"{{{{clicks_category_{idx}}}}}"] = f"{categoryRow.get('clicks', 0):,}"
            replacements[f"{{{{spend_category_{idx}}}}}"] = f"${int(categoryRow.get('spend', 0)):,}"
            replacements[f"{{{{sales_category_{idx}}}}}"] = f"${int(categoryRow.get('sales', 0)):,}"
            replacements[f"{{{{units_category_{idx}}}}}"] = f"{categoryRow.get('units', 0):,}"
            replacements[f"{{{{ROAS_category_{idx}}}}}"] = f"${categoryRow.get('ROAS Demand', 0):,.2f}"
            replacements[f"{{{{CPC_category_{idx}}}}}"] = f"${categoryRow.get('CPC Demand', 0):,.2f}"
            

    slide = prs.slides[slide_index]
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_text_in_shape(shape, replacements)
        
        if shape.has_table:
            replace_text_in_table(shape.table, replacements)


    import re
    # Remove any table rows that still have a {{…}} token
    for shape in slide.shapes:
        if not shape.has_table:
            continue
    
        tbl = shape.table
        # make a static list of (row, xml) so we can remove safely
        rows_with_xml = [(row, row._tr) for row in tbl.rows]
    
        for row, row_tr in rows_with_xml:
            # collect all the text in the row
            texts = []
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    texts.append(para.text)
    
            # if any cell in that row has a {{…}} placeholder
            if any(re.search(r"\{\{[^}]+\}\}", text) for text in texts):
                # remove the <w:tr> element from the table
                tbl._tbl.remove(row_tr)


    #print(replacements)
    print("slide", slide_index + 1, "done")

#Formats capout times into 12hr times.
def format_hour_12h(hour_24):
    """
    hour_24: int or numeric string (0–23)
    returns: e.g. "8:00 PM", "12:00 AM", "11:00 AM"
    """
    h = int(hour_24) % 24
    suffix = "AM" if h < 12 else "PM"
    h12 = h % 12 or 12
    return f"{h12}:00 {suffix}"


#converts dates into human times.
from datetime import datetime
def ordinal(n: int) -> str:
    # returns e.g. "1st", "2nd", "3rd", "4th", …, "11th", "21st", etc.
    if 10 <= (n % 100) <= 20:
        suffix = "th"
    else:
        suffix = {1: "st", 2: "nd", 3: "rd"}.get(n % 10, "th")
    return f"{n}{suffix}"

def human_date(datestr: str) -> str:
    """
    Convert "YYYY-MM-DD" → "Month D<ordinal>, YYYY"
    e.g. "2025-01-01" → "January 1st, 2025"
    """
    dt = datetime.strptime(datestr, "%Y-%m-%d")
    month = dt.strftime("%B")         # full month name
    day_ord = ordinal(dt.day)         # e.g. "1st"
    year = dt.year
    return f"{month} {day_ord}, {year}"
    
# Main execution
prs = Presentation("template.pptx")

# Slide 1: YTD Campaign Snapshot
fill_slide(prs, 0, "2025-01-01", "2025-05-06", use_line_items=False, by_keyword=False, by_category=False)

# Slide 2 - Month Snapshot
fill_slide(prs, 1, "2025-04-01", "2025-04-30", use_line_items=False, by_keyword=False, by_category=False)

# Slide 2: line-item-level
fill_slide(prs, 2, "2025-04-01", "2025-05-06", use_line_items=True, by_keyword=False, by_category=False)

#Slide 3: keyword level
fill_slide(prs, 3, "2025-04-01", "2025-05-06", use_line_items=False, by_keyword=True, by_category=False)

#Slide 4: category level
fill_slide(prs, 4, "2025-04-01", "2025-05-06", use_line_items=False, by_keyword=False, by_category=True)

prs.save(presentationName)
print("Presentation updated successfully.")
